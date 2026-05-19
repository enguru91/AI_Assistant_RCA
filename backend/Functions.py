# Author: Eshan Gurusinghe
# Email: engurusinghe91@gmail.com
# Version: 3.0 — Migrated from Ericsson ELI API to local Ollama stack -> Date: 2026-05-17
# Version: 3.1 — Dynamic model selector and model installation/uninstallation
#              — Added disk space validation before installation -> Date: 2026-05-18
# Version: 3.2 — Extended upload support: txt, csv, md, xml, json, pptx, xlsx, docx, pdf
#              — Added file validation layer: per-file 200 MB cap, 400 MB batch cap
#              — Added magic-bytes content sniffing for binary formats
#              — Added extract_text_from_file() dispatcher (all 9 formats)
#              — save_uploaded_files() now extracts text; raw bytes are never stored
#              — check_for_uploaded_files() returns per-file extraction warnings
#              — Updated delete_remaining_resources() to cover new resource extensions
#              -> Date: 2026-05-18
#
# New runtime dependencies (add to requirements.txt):
#   pypdf>=4.0          — PDF text extraction
#   python-docx>=1.1    — DOCX text extraction
#   python-pptx>=0.6    — PPTX text extraction
#   openpyxl>=3.1       — XLSX text extraction
#   (sentence-transformers, openai, streamlit, python-dotenv — unchanged)

import os
import io
import csv
import json
import pickle
import glob
import hashlib
import subprocess
import shutil
import numpy as np
import streamlit as st
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
from sentence_transformers import SentenceTransformer, CrossEncoder
import xml.etree.ElementTree as ET

# ---------------------------------------------------------------------------
# Environment & path constants
# ---------------------------------------------------------------------------

load_dotenv()

resources_path   = "./Include/ProcedureFiles/"
embeddings_path  = "./Include/Outputs/Embeddings/"
output_dir       = "./Include/Outputs/Chats/"
current_datetime = datetime.now().strftime("%Y%m%d_%H%M%S")
filename_prefix  = "ResourcesFile_"  + current_datetime
checksum_file    = "current_checksum_" + current_datetime

# Ensure output directories exist on startup
os.makedirs(resources_path,  exist_ok=True)
os.makedirs(embeddings_path, exist_ok=True)
os.makedirs(output_dir,      exist_ok=True)

# ---------------------------------------------------------------------------
# File upload — validation constants
# ---------------------------------------------------------------------------

# Hard size limits
MAX_SINGLE_FILE_MB: int = 100
MAX_BATCH_MB:       int = 200
MAX_SINGLE_FILE_B:  int = MAX_SINGLE_FILE_MB * 1024 * 1024
MAX_BATCH_B:        int = MAX_BATCH_MB        * 1024 * 1024

# Canonical set of accepted extensions (lower-case, no leading dot)
ALLOWED_EXTENSIONS: frozenset[str] = frozenset({
    "txt", "csv", "md",
    "xml", "json",
    "pptx",                 # OOXML PowerPoint — legacy .ppt is not supported
    "xlsx",                 # OOXML Excel
    "docx",                 # OOXML Word
    "pdf",
})

# Magic-byte signatures used for content-sniffing of binary formats.
# Keys are byte prefixes (first 8 bytes); values are the expected extension group.
_MAGIC_BYTES: dict[bytes, str] = {
    b"%PDF":      "pdf",
    b"PK\x03\x04": "ooxml",   # All three OOXML formats (docx/xlsx/pptx) are ZIP archives
}

_OOXML_EXTENSIONS: frozenset[str] = frozenset({"docx", "xlsx", "pptx"})

# ---------------------------------------------------------------------------
# File upload — validation
# ---------------------------------------------------------------------------

def validate_uploaded_files(uploaded_files) -> list[str]:
    """
    Validate a list of Streamlit UploadedFile objects before any processing.

    Checks performed (in order):
      1. Extension is in ALLOWED_EXTENSIONS.
      2. Individual file size does not exceed MAX_SINGLE_FILE_MB.
      3. Magic bytes match the declared extension (binary formats only).
      4. Aggregate batch size does not exceed MAX_BATCH_MB.

    Parameters
    ----------
    uploaded_files : list[UploadedFile]
        Files returned by st.file_uploader().

    Returns
    -------
    list[str]
        Human-readable error strings. An empty list means all files are valid.
    """
    errors: list[str] = []
    total_bytes: int = 0

    for f in uploaded_files:
        name = f.name
        ext  = os.path.splitext(name)[1].lstrip(".").lower()
        size = f.size   # Streamlit always exposes .size

        # ── 1. Extension whitelist ────────────────────────────────────────
        if ext not in ALLOWED_EXTENSIONS:
            errors.append(
                f"❌ '{name}': unsupported file type '.{ext}'. "
                f"Accepted types: {', '.join(sorted(ALLOWED_EXTENSIONS))}."
            )
            continue  # No point in further checks for this file

        # ── 2. Individual size cap ────────────────────────────────────────
        if size > MAX_SINGLE_FILE_B:
            size_mb = size / (1024 ** 2)
            errors.append(
                f"❌ '{name}': file size {size_mb:.1f} MB exceeds the "
                f"{MAX_SINGLE_FILE_MB} MB per-file limit."
            )

        # ── 3. Magic-byte content sniffing (binary formats only) ──────────
        if ext in _OOXML_EXTENSIONS or ext == "pdf":
            header = f.read(8)
            f.seek(0)   # Always reset the stream pointer after peeking

            if ext == "pdf":
                if not header.startswith(b"%PDF"):
                    errors.append(
                        f"❌ '{name}': file does not start with a valid PDF signature. "
                        "It may be corrupt or misnamed."
                    )
            elif ext in _OOXML_EXTENSIONS:
                if not header.startswith(b"PK\x03\x04"):
                    errors.append(
                        f"❌ '{name}': file does not appear to be a valid Office Open XML "
                        f"file (expected ZIP/OOXML signature). It may be corrupt, "
                        f"or a legacy binary format (e.g. .ppt / .xls / .doc) renamed "
                        f"with a modern extension. Please re-save as .{ext} first."
                    )
        else:
            # For text-based formats reset is not needed, but do it for safety
            f.seek(0)

        total_bytes += size

    # ── 4. Batch size cap ─────────────────────────────────────────────────
    if total_bytes > MAX_BATCH_B:
        total_mb = total_bytes / (1024 ** 2)
        errors.append(
            f"❌ Total batch size {total_mb:.1f} MB exceeds the {MAX_BATCH_MB} MB "
            f"batch limit. Remove some files and try again."
        )

    return errors


# ---------------------------------------------------------------------------
# File upload — text extraction
# ---------------------------------------------------------------------------

def extract_text_from_file(uploaded_file) -> str:
    """
    Extract all readable text from an UploadedFile object.

    Supported formats
    -----------------
    txt, md  — UTF-8 decode (BOM-aware, errors replaced)
    csv      — Rows joined as comma-separated strings
    json     — Pretty-printed JSON tree (validates structure)
    xml      — All element text nodes concatenated
    xlsx     — All sheets, rows as tab-separated strings
    docx     — Paragraphs + table cells
    pptx     — All slide text shapes, labelled by slide number
    pdf      — Page text via pypdf (text-layer PDFs only)

    Parameters
    ----------
    uploaded_file : UploadedFile
        A Streamlit UploadedFile whose stream is positioned at the start.

    Returns
    -------
    str
        Extracted text content.

    Raises
    ------
    ValueError
        If the file content cannot be parsed or a required library is missing.
    """
    name = uploaded_file.name
    ext  = os.path.splitext(name)[1].lstrip(".").lower()
    raw  = uploaded_file.read()
    uploaded_file.seek(0)  # Reset so callers can re-read if needed

    # ── Plain text ─────────────────────────────────────────────────────────
    if ext in ("txt", "md"):
        # Strip UTF-8 BOM if present
        return raw.decode("utf-8-sig", errors="replace")

    # ── CSV ────────────────────────────────────────────────────────────────
    if ext == "csv":
        text_io = io.StringIO(raw.decode("utf-8-sig", errors="replace"))
        reader  = csv.reader(text_io)
        rows    = [", ".join(cell.strip() for cell in row) for row in reader if any(row)]
        return "\n".join(rows)

    # ── JSON ───────────────────────────────────────────────────────────────
    if ext == "json":
        try:
            data = json.loads(raw.decode("utf-8-sig", errors="replace"))
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError as exc:
            raise ValueError(f"Invalid JSON in '{name}': {exc}") from exc

    # ── XML ────────────────────────────────────────────────────────────────
    if ext == "xml":
        try:
            root = ET.fromstring(raw)
            parts: list[str] = []
            for elem in root.iter():
                if elem.text  and elem.text.strip():
                    parts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    parts.append(elem.tail.strip())
            if not parts:
                raise ValueError(f"No text content found in XML file '{name}'.")
            return "\n".join(parts)
        except ET.ParseError as exc:
            raise ValueError(f"Malformed XML in '{name}': {exc}") from exc

    # ── PDF ────────────────────────────────────────────────────────────────
    if ext == "pdf":
        try:
            import pypdf
        except ImportError as exc:
            raise ValueError(
                "pypdf is not installed. Run: pip install 'pypdf>=4.0'"
            ) from exc
        try:
            reader = pypdf.PdfReader(io.BytesIO(raw))
            if reader.is_encrypted:
                raise ValueError(
                    f"PDF '{name}' is password-protected. "
                    "Please decrypt it before uploading."
                )
            pages = [page.extract_text() or "" for page in reader.pages]
            text  = "\n\n".join(p.strip() for p in pages if p.strip())
            if not text:
                raise ValueError(
                    f"No extractable text found in '{name}'. "
                    "The PDF may be image-only (scanned). "
                    "Please use a PDF with a text layer, or run OCR first."
                )
            return text
        except pypdf.errors.PdfReadError as exc:
            raise ValueError(f"Could not read PDF '{name}': {exc}") from exc

    # ── DOCX ───────────────────────────────────────────────────────────────
    if ext == "docx":
        try:
            import docx as python_docx
        except ImportError as exc:
            raise ValueError(
                "python-docx is not installed. Run: pip install 'python-docx>=1.1'"
            ) from exc
        try:
            doc   = python_docx.Document(io.BytesIO(raw))
            parts = []

            # Body paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())

            # Tables (each row as pipe-delimited)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.strip():
                        parts.append(row_text)

            if not parts:
                raise ValueError(f"No readable text found in DOCX '{name}'.")
            return "\n".join(parts)
        except Exception as exc:
            raise ValueError(f"Could not read DOCX '{name}': {exc}") from exc

    # ── XLSX ───────────────────────────────────────────────────────────────
    if ext == "xlsx":
        try:
            import openpyxl
        except ImportError as exc:
            raise ValueError(
                "openpyxl is not installed. Run: pip install 'openpyxl>=3.1'"
            ) from exc
        try:
            wb    = openpyxl.load_workbook(
                io.BytesIO(raw), read_only=True, data_only=True
            )
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()

            if not parts:
                raise ValueError(f"No data found in XLSX '{name}'.")
            return "\n".join(parts)
        except Exception as exc:
            raise ValueError(f"Could not read XLSX '{name}': {exc}") from exc

    # ── PPTX ───────────────────────────────────────────────────────────────
    if ext == "pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ValueError(
                "python-pptx is not installed. Run: pip install 'python-pptx>=0.6'"
            ) from exc
        try:
            prs    = Presentation(io.BytesIO(raw))
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    # Text frames (title, body, text boxes)
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs).strip()
                            if line:
                                texts.append(line)
                    # Tables inside slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip():
                                texts.append(row_text)
                if texts:
                    slides.append(f"[Slide {idx}]\n" + "\n".join(texts))

            if not slides:
                raise ValueError(f"No readable text found in PPTX '{name}'.")
            return "\n\n".join(slides)
        except Exception as exc:
            raise ValueError(f"Could not read PPTX '{name}': {exc}") from exc

    # ── Fallback (should never reach here after extension validation) ──────
    raise ValueError(
        f"No extractor available for '.{ext}'. "
        f"Supported types: {', '.join(sorted(ALLOWED_EXTENSIONS))}."
    )


# ---------------------------------------------------------------------------
# Local model initialisation
# ---------------------------------------------------------------------------

# Embedding model — converts text chunks into 384-dimensional vectors.
# Auto-downloads (~130 MB) on first use, then cached locally.
embed_model = SentenceTransformer("BAAI/bge-small-en-v1.5")

# Cross-encoder reranking model — scores (query, document) pairs for relevance.
# Auto-downloads (~80 MB) on first use, then cached locally.
rerank_model = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")

# ---------------------------------------------------------------------------
# Ollama client initialisation (OpenAI-compatible)
# ---------------------------------------------------------------------------

# Ollama exposes an OpenAI-compatible API at /v1.
# LLM_API_KEY is required by the SDK but ignored by Ollama — no real auth.
ollama_client = OpenAI(
    base_url=os.environ.get("LLM_BASE_URL", "http://localhost:11434/v1"),
    api_key=os.environ.get("LLM_API_KEY",   "ollama"),
)

# ---------------------------------------------------------------------------
# Ollama health check
# ---------------------------------------------------------------------------

def check_ollama_connection():
    """
    Check whether the Ollama server is reachable and return available models.
    Returns (True, [model_names]) on success, or (False, error_message) on failure.
    """
    try:
        models      = ollama_client.models.list()
        model_names = [m.id for m in models.data]
        print(f"Ollama connected. Available models: {model_names}")
        return True, model_names
    except Exception as ex:
        print(f"Ollama connection failed: {ex}")
        return False, str(ex)

# ---------------------------------------------------------------------------
# Ollama model catalogue
# ---------------------------------------------------------------------------

# A curated list of popular, publicly available Ollama models.
# Format: {"display_label": "ollama_pull_name"}
# Models the user already has installed are filtered out dynamically.
KNOWN_OLLAMA_MODELS = {
    # ── Small & fast (good for CPU-only machines) ──────────────────────────
    "LLaMA 3.1 8B  (~4.7 GB)":    "llama3.1:8b",
    "LLaMA 3.2 3B  (~2.0 GB)":    "llama3.2:3b",
    "LLaMA 3.2 1B  (~1.3 GB)":    "llama3.2:1b",
    "Mistral 7B    (~4.1 GB)":    "mistral:latest",
    "Gemma 2 2B    (~1.6 GB)":    "gemma2:2b",
    "Phi-3 Mini    (~2.2 GB)":    "phi3:mini",
    "Qwen 2.5 7B   (~4.4 GB)":    "qwen2.5:7b",
    "DeepSeek-R1 7B (~4.7 GB)":   "deepseek-r1:7b",
    # ── Medium (GPU recommended) ──────────────────────────────────────────
    "LLaMA 3.1 70B (~40 GB)":     "llama3.1:70b",
    "Mistral Nemo 12B (~7.2 GB)": "mistral-nemo:latest",
    "Gemma 2 9B    (~5.5 GB)":    "gemma2:9b",
    "Phi-4 14B     (~9.1 GB)":    "phi4:latest",
    "Qwen 2.5 14B  (~8.9 GB)":    "qwen2.5:14b",
    "DeepSeek-R1 14B (~9.0 GB)":  "deepseek-r1:14b",
    # ── Large (high-end GPU) ──────────────────────────────────────────────
    "Qwen 2.5 32B  (~20 GB)":     "qwen2.5:32b",
    "DeepSeek-R1 32B (~20 GB)":   "deepseek-r1:32b",
    "LLaMA 3.3 70B (~43 GB)":     "llama3.3:70b",
}

# Approximate download sizes in GB for each model in the catalogue above.
# Used by check_disk_space() to validate free space before starting a download.
MODEL_SIZES_GB = {
    "llama3.1:8b":          4.7,
    "llama3.2:3b":          2.0,
    "llama3.2:1b":          1.3,
    "mistral:latest":       4.1,
    "gemma2:2b":            1.6,
    "phi3:mini":            2.2,
    "qwen2.5:7b":           4.4,
    "deepseek-r1:7b":       4.7,
    "llama3.1:70b":        40.0,
    "mistral-nemo:latest":  7.2,
    "gemma2:9b":            5.5,
    "phi4:latest":          9.1,
    "qwen2.5:14b":          8.9,
    "deepseek-r1:14b":      9.0,
    "qwen2.5:32b":         20.0,
    "deepseek-r1:32b":     20.0,
    "llama3.3:70b":        43.0,
}

# ---------------------------------------------------------------------------
# Dynamic Ollama model listing
# ---------------------------------------------------------------------------

def list_ollama_models() -> list[str]:
    """
    Fetch the names of all models currently installed in Ollama.

    Returns
    -------
    list[str]
        Sorted list of installed model name strings (e.g. ["llama3.1:8b", "mistral:latest"]).
        Returns an empty list if Ollama is unreachable or no models are installed.
    """
    try:
        models      = ollama_client.models.list()
        model_names = sorted([m.id for m in models.data])
        print(f"Installed Ollama models: {model_names}")
        return model_names
    except Exception as ex:
        print(f"Could not list Ollama models: {ex}")
        return []

# ---------------------------------------------------------------------------
# Ollama model installation
# ---------------------------------------------------------------------------

def install_ollama_model(model_name: str) -> tuple[bool, str]:
    """
    Pull (download and install) an Ollama model by running:
        ollama pull <model_name>

    This is a blocking call — it waits until the download completes.

    Returns
    -------
    (True, success_message)   on success
    (False, error_message)    on failure
    """
    print(f"Starting install: ollama pull {model_name}")
    try:
        result = subprocess.run(
            ["ollama", "pull", model_name],
            capture_output=True,
            text=True,
            timeout=3600,   # 1-hour hard timeout — large models can be slow
        )
        if result.returncode == 0:
            msg = f"Model '{model_name}' installed successfully."
            print(msg)
            return True, msg
        else:
            err = result.stderr.strip() or result.stdout.strip()
            msg = f"ollama pull failed for '{model_name}': {err}"
            print(msg)
            return False, msg
    except FileNotFoundError:
        msg = "ollama command not found. Is Ollama installed and on your PATH?"
        print(msg)
        return False, msg
    except subprocess.TimeoutExpired:
        msg = f"Install timed out after 1 hour for model '{model_name}'."
        print(msg)
        return False, msg
    except Exception as ex:
        msg = f"Unexpected error during install of '{model_name}': {ex}"
        print(msg)
        return False, msg

# ---------------------------------------------------------------------------
# Disk space validation
# ---------------------------------------------------------------------------

def check_disk_space(model_name: str) -> dict:
    """
    Validate whether the target drive has enough free space to install a model.

    Logic
    -----
    1. Resolve the drive to check from the OLLAMA_MODELS environment variable.
       Falls back to the user home directory if the variable is not set or the
       path does not exist yet.
    2. Look up the required GB from MODEL_SIZES_GB. If the model is not in the
       catalogue (e.g. a custom name), treat required size as 0 and warn.
    3. Three possible outcomes:
       - BLOCKED  (can_install=False): free space < required size.
       - WARNING  (can_install=True,  warning=True): fits but ≤ 10% would remain.
       - OK       (can_install=True,  warning=False): plenty of space.

    Returns
    -------
    dict with keys: can_install, warning, message, required_gb, free_gb,
                    total_gb, remaining_pct
    """
    # ── 1. Resolve the drive path ─────────────────────────────────────────
    ollama_models_env = os.environ.get("OLLAMA_MODELS", "")
    if ollama_models_env and os.path.exists(ollama_models_env):
        check_path = ollama_models_env
    else:
        default_path = os.path.join(os.path.expanduser("~"), ".ollama", "models")
        check_path   = default_path if os.path.exists(default_path) \
                       else os.path.expanduser("~")

    print(f"Disk space check path: {check_path}")

    # ── 2. Get drive stats ─────────────────────────────────────────────────
    try:
        usage    = shutil.disk_usage(check_path)
        total_gb = usage.total / (1024 ** 3)
        free_gb  = usage.free  / (1024 ** 3)
    except Exception as ex:
        print(f"Could not read disk usage for {check_path}: {ex}")
        return {
            "can_install":   True,
            "warning":       True,
            "message":       f"Could not read disk space ({ex}). Verify manually before installing.",
            "required_gb":   0.0,
            "free_gb":       0.0,
            "total_gb":      0.0,
            "remaining_pct": 0.0,
        }

    # ── 3. Look up required size ───────────────────────────────────────────
    required_gb = MODEL_SIZES_GB.get(model_name)
    if required_gb is None:
        print(f"Model '{model_name}' not found in MODEL_SIZES_GB — size unknown.")
        remaining_pct = (free_gb / total_gb * 100) if total_gb > 0 else 0.0
        return {
            "can_install":   True,
            "warning":       True,
            "message":       (
                f"Model size for '{model_name}' is unknown — it is not in the built-in catalogue. "
                f"Current free space: {free_gb:.1f} GB. Verify the size manually before proceeding."
            ),
            "required_gb":   0.0,
            "free_gb":       free_gb,
            "total_gb":      total_gb,
            "remaining_pct": remaining_pct,
        }

    # ── 4. Evaluate space ──────────────────────────────────────────────────
    remaining_after_gb  = free_gb - required_gb
    remaining_after_pct = (remaining_after_gb / total_gb * 100) if total_gb > 0 else 0.0

    print(
        f"Disk check — required: {required_gb:.1f} GB | free: {free_gb:.1f} GB | "
        f"remaining after: {remaining_after_gb:.1f} GB ({remaining_after_pct:.1f}%)"
    )

    if free_gb < required_gb:
        return {
            "can_install":   False,
            "warning":       False,
            "message":       (
                f"Not enough disk space to install '{model_name}'. "
                f"Required: {required_gb:.1f} GB — Available: {free_gb:.1f} GB — "
                f"Shortfall: {required_gb - free_gb:.1f} GB. "
                f"Free up space on the drive and try again."
            ),
            "required_gb":   required_gb,
            "free_gb":       free_gb,
            "total_gb":      total_gb,
            "remaining_pct": remaining_after_pct,
        }

    if remaining_after_pct <= 10.0:
        return {
            "can_install":   True,
            "warning":       True,
            "message":       (
                f"Low disk space warning. After installing '{model_name}' "
                f"({required_gb:.1f} GB), only {remaining_after_pct:.1f}% "
                f"({remaining_after_gb:.1f} GB) of the drive will remain free. "
                f"Running with less than 10% free space may affect system performance."
            ),
            "required_gb":   required_gb,
            "free_gb":       free_gb,
            "total_gb":      total_gb,
            "remaining_pct": remaining_after_pct,
        }

    return {
        "can_install":   True,
        "warning":       False,
        "message":       (
            f"Sufficient disk space available. "
            f"Required: {required_gb:.1f} GB — Free: {free_gb:.1f} GB — "
            f"Remaining after install: {remaining_after_pct:.1f}%."
        ),
        "required_gb":   required_gb,
        "free_gb":       free_gb,
        "total_gb":      total_gb,
        "remaining_pct": remaining_after_pct,
    }

# ---------------------------------------------------------------------------
# Ollama model uninstall
# ---------------------------------------------------------------------------

def uninstall_ollama_model(model_name: str) -> tuple[bool, str]:
    """
    Remove an installed Ollama model by running:
        ollama rm <model_name>

    Returns
    -------
    (True, success_message)   on success
    (False, error_message)    on failure
    """
    print(f"Uninstalling model: ollama rm {model_name}")
    try:
        result = subprocess.run(
            ["ollama", "rm", model_name],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode == 0:
            msg = f"Model '{model_name}' uninstalled successfully."
            print(msg)
            return True, msg
        else:
            err = result.stderr.strip() or result.stdout.strip()
            msg = f"ollama rm failed for '{model_name}': {err}"
            print(msg)
            return False, msg
    except FileNotFoundError:
        msg = "ollama command not found. Is Ollama installed and on your PATH?"
        print(msg)
        return False, msg
    except subprocess.TimeoutExpired:
        msg = f"Uninstall timed out for model '{model_name}'."
        print(msg)
        return False, msg
    except Exception as ex:
        msg = f"Unexpected error during uninstall of '{model_name}': {ex}"
        print(msg)
        return False, msg

# ---------------------------------------------------------------------------
# File & resource management
# ---------------------------------------------------------------------------

def delete_remaining_resources():
    """
    Delete leftover ResourcesFile_, current_checksum_, and current_embeddings_
    files from previous sessions. Called once at app startup.
    """
    for filename in os.listdir(resources_path):
            file_path = os.path.join(resources_path, filename)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.remove(file_path)
                    print(f"Deleted file: {file_path}")
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
                    print(f"Deleted directory: {file_path}")
            except Exception as ex:
                print(f"Failed to delete {file_path}: {ex}")
    else:
         print(f"Directory does not exist, skipping cleanup: {resources_path}")

    cleanup_targets = [
        (output_dir,      "chat_output_",        ".txt"),
        (embeddings_path, "current_embeddings_", ".pkl"),
    ]
    
    for directory, prefix, extension in cleanup_targets:
        if not os.path.isdir(directory):
            print(f"Directory does not exist, skipping cleanup: {directory}")
            continue
        for filename in os.listdir(directory):
            # Using endswith might be too restrictive if you just want prefix matching
            # but keeping original logic here.
            if filename.startswith(prefix) and filename.endswith(extension):
                file_path = os.path.join(directory, filename)
                try:
                    os.remove(file_path)
                    print(f"Deleted targeted file: {file_path}")
                except Exception as ex:
                    print(f"Failed to delete {file_path}: {ex}")

def save_uploaded_files(uploaded_files) -> tuple[str | None, list[str]]:
    """
    Extract text from each uploaded file and persist it as a ResourcesFile_ .txt.

    All file types (pdf, docx, xlsx, pptx, csv, xml, json, md, txt) are converted
    to plain text at this stage. Downstream functions (embeddings, reranker) only
    ever see UTF-8 .txt files regardless of the original format.

    Each saved file is prefixed with a '[Source: <filename>]' header so the LLM
    knows which document a text chunk originates from.

    Parameters
    ----------
    uploaded_files : list[UploadedFile]
        Pre-validated Streamlit file objects.

    Returns
    -------
    (checksum, warnings)
        checksum : str | None
            MD5 hex digest of all saved resource files, or None if none were saved.
        warnings : list[str]
            Per-file extraction warnings (e.g. a file that could not be parsed).
            Does not include validation errors — those are checked earlier.
    """
    warnings: list[str] = []

    # Remove any existing ResourcesFile_ entries before saving new ones
    for file in glob.glob(os.path.join(resources_path, "ResourcesFile_*")):
        try:
            os.remove(file)
            print(f"Deleted old resource file: {file}")
        except Exception as ex:
            print(f"Error deleting resource file {file}: {ex}")

    saved_count = 0
    for index, uploaded_file in enumerate(uploaded_files):
        source_name = uploaded_file.name

        # ── Extract text from the file ────────────────────────────────────
        try:
            text_content = extract_text_from_file(uploaded_file)
        except ValueError as exc:
            msg = f"⚠️ '{source_name}': skipped — {exc}"
            print(msg)
            warnings.append(msg)
            continue    # Do not save a resource file for this upload

        if not text_content.strip():
            msg = f"⚠️ '{source_name}': skipped — no text content was extracted."
            print(msg)
            warnings.append(msg)
            continue

        # ── Persist as UTF-8 .txt with source header ──────────────────────
        new_filename = f"{filename_prefix}_{index + 1}.txt"
        file_path    = os.path.join(resources_path, new_filename)
        try:
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(f"[Source: {source_name}]\n\n{text_content}\n")
            print(f"Saved extracted text ({len(text_content):,} chars): {file_path}")
            saved_count += 1
        except OSError as exc:
            msg = f"⚠️ '{source_name}': could not write resource file — {exc}"
            print(msg)
            warnings.append(msg)

    if saved_count == 0:
        print("No resource files were saved.")
        return None, warnings

    # ── Compute checksum across all newly saved files ─────────────────────
    new_checksum = compute_resources_checksum(resources_path)
    if not new_checksum:
        print("Error: checksum computation returned empty.")
        return None, warnings

    # Write checksum to any existing current_checksum_ file
    for filename in sorted(os.listdir(resources_path)):
        if filename.startswith("current_checksum_") and filename.endswith(".txt"):
            filepath = os.path.join(resources_path, filename)
            try:
                with open(filepath, "w") as f:
                    f.write(new_checksum)
                print(f"Checksum written to {filepath}: {new_checksum}")
            except Exception as ex:
                print(f"Error writing checksum to {filepath}: {ex}")

    return new_checksum, warnings

def compute_resources_checksum(directory: str) -> str | None:
    """
    Compute an MD5 hash across all ResourcesFile_ .txt files in the directory.
    Returns the hex digest string, or None if no files are found.
    """
    hash_md5  = hashlib.md5()
    found_any = False
    for filename in sorted(os.listdir(directory)):
        if filename.startswith("ResourcesFile_") and filename.endswith(".txt"):
            filepath  = os.path.join(directory, filename)
            found_any = True
            with open(filepath, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
    if not found_any:
        print("No ResourcesFile_ files found for checksum computation.")
        return None
    return hash_md5.hexdigest()

def check_exists_checksum_file(directory: str) -> str | None:
    """
    Find the current checksum file in the given directory.
    If none exists, create a new empty one.
    Returns the full path to the checksum file.
    """
    try:
        matching_files = [
            f for f in os.listdir(directory)
            if f.startswith("current_checksum_") and f.endswith(".txt")
        ]
        if matching_files:
            matching_files.sort()
            return os.path.join(directory, matching_files[0])

        print("No checksum file found. Creating a new one.")
        new_checksum_filepath = os.path.join(
            directory, f"current_checksum_{current_datetime}.txt"
        )
        with open(new_checksum_filepath, "w") as f:
            f.write("")
        print(f"Created new empty checksum file: {new_checksum_filepath}")
        return new_checksum_filepath

    except FileNotFoundError:
        print(f"Directory does not exist: {directory}")
        return None
    except Exception as ex:
        print(f"Error searching for checksum file: {ex}")
        return None


def check_for_uploaded_files(uploaded_files) -> tuple | None:
    """
    Orchestrate the full upload flow:
      1. Save (extract + persist) uploaded files to disk.
      2. Compare new checksum against the stored one.
      3. Return cached embeddings if unchanged, or create new ones if changed.

    Parameters
    ----------
    uploaded_files : list[UploadedFile]
        Pre-validated Streamlit file objects (validation has already been run
        in the UI layer via validate_uploaded_files() before this is called).

    Returns
    -------
    (texts, vectors, warnings) — on success
    None                       — if no files were provided
    """
    current_checksum_path = check_exists_checksum_file(resources_path)

    if not uploaded_files:
        print("No uploaded files detected.")
        return None

    print("Uploaded files detected.")

    # Read the stored checksum (if any)
    current_checksum_value = None
    if current_checksum_path and os.path.exists(current_checksum_path):
        try:
            with open(current_checksum_path, "r") as f:
                current_checksum_value = f.read().strip()
            print(f"Stored checksum: {current_checksum_value}")
        except Exception as ex:
            print(f"Error reading checksum file: {ex}")

    # Save files (text extraction happens here) and compute new checksum
    new_checksum_value, extraction_warnings = save_uploaded_files(uploaded_files)
    print(f"New checksum: {new_checksum_value}")

    if new_checksum_value is None:
        # All files failed extraction — nothing to embed
        return [], np.array([]), extraction_warnings

    if current_checksum_value and current_checksum_value == new_checksum_value:
        print("Checksum matches — loading cached embeddings.")
        texts, vectors = load_embeddings_from_file()
        return texts, vectors, extraction_warnings

    print("Checksum differs — creating new embeddings.")
    embeddings_data = create_embeddings()

    # Persist the new checksum
    if current_checksum_path:
        try:
            with open(current_checksum_path, "w") as f:
                f.write(new_checksum_value)
            print(f"Updated stored checksum to: {new_checksum_value}")
        except Exception as ex:
            print(f"Error updating checksum file: {ex}")

    texts, vectors = embeddings_data
    return texts, vectors, extraction_warnings


def load_resources(directory: str) -> list[str]:
    """
    Load all ResourcesFile_ .txt files from the given directory.
    Returns a list of strings (one per file).
    """
    texts = []
    for filename in os.listdir(directory):
        if filename.startswith("ResourcesFile_") and filename.endswith(".txt"):
            file_path = os.path.join(directory, filename)
            with open(file_path, "r", encoding="utf-8") as f:
                texts.append(f.read())
    print(f"Loaded {len(texts)} resource file(s).")
    return texts

# ---------------------------------------------------------------------------
# Embedding — local sentence-transformers model
# ---------------------------------------------------------------------------

def create_embeddings_local(texts: list[str]) -> dict:
    """
    Encode a list of text strings into embedding vectors using the local
    BAAI/bge-small-en-v1.5 model. Runs entirely on CPU with no API call.
    Returns a dict {"vectors": list_of_lists} to match the original interface.
    """
    try:
        vectors = embed_model.encode(texts, convert_to_numpy=True)
        print(f"Created embeddings for {len(texts)} text(s). Shape: {vectors.shape}")
        return {"vectors": vectors.tolist()}
    except Exception as ex:
        print(f"Error during embedding: {ex}")
        return {}


def create_embeddings(embeddings_dir: str = embeddings_path) -> tuple:
    """
    Delete any existing embedding .pkl files, create fresh embeddings from
    all resource files, and save them to a new timestamped .pkl file.
    Returns (texts, numpy_vectors) tuple.
    """
    for file in glob.glob(os.path.join(embeddings_dir, "current_embeddings_*")):
        try:
            os.remove(file)
            print(f"Deleted old embedding file: {file}")
        except Exception as ex:
            print(f"Error deleting embedding file {file}: {ex}")

    print("Creating new embeddings...")
    texts = load_resources(resources_path)
    if not texts:
        print("No resource texts found. Cannot create embeddings.")
        return [], np.array([])

    embeddings = create_embeddings_local(texts)

    if not embeddings.get("vectors"):
        print("Embedding creation failed.")
        return [], np.array([])

    new_filename  = f"current_embeddings_{current_datetime}.pkl"
    new_file_path = os.path.join(embeddings_dir, new_filename)
    with open(new_file_path, "wb") as f:
        pickle.dump({"texts": texts, "vectors": embeddings["vectors"]}, f)
    print(f"Saved new embeddings to: {new_filename}")
    return texts, np.array(embeddings["vectors"])


@st.cache_data
def load_embeddings_from_file(embeddings_dir: str = embeddings_path) -> tuple:
    """
    Load embeddings from the most recent .pkl file in embeddings_dir.
    Falls back to creating new embeddings if no .pkl file is found.
    Result is cached by Streamlit to avoid reloading on every rerun.
    """
    pkl_files = [
        f for f in os.listdir(embeddings_dir)
        if f.startswith("current_embeddings_") and f.endswith(".pkl")
    ]
    if pkl_files:
        pkl_files.sort()
        embeddings_file = os.path.join(embeddings_dir, pkl_files[-1])
        with open(embeddings_file, "rb") as f:
            data = pickle.load(f)
        print(f"Loaded cached embeddings from: {embeddings_file}")
        return data["texts"], np.array(data["vectors"])

    print("No cached embeddings found. Creating new ones.")
    return create_embeddings(embeddings_dir)

# ---------------------------------------------------------------------------
# Reranking — local CrossEncoder model
# ---------------------------------------------------------------------------

def local_score_rank(query: str, texts: list[str]) -> dict:
    """
    Score each (query, document) pair using the local CrossEncoder model.
    Returns a dict {"scores": list_of_floats}.
    Runs entirely on CPU with no API call.
    """
    try:
        pairs  = [(query, doc) for doc in texts]
        scores = rerank_model.predict(pairs)
        print(f"Reranking complete. Scores: {scores}")
        return {"scores": scores.tolist()}
    except Exception as ex:
        print(f"Error during reranking: {ex}")
        return {}


def get_most_relevant_docs(query: str, texts: list[str]) -> list[str]:
    """
    Return the uploaded documents sorted by relevance to the query,
    most relevant first, using the local CrossEncoder reranker.
    """
    if not texts:
        print("No texts available for reranking.")
        return []

    reranked_response = local_score_rank(query=query, texts=texts)
    if not reranked_response:
        return []

    scores = reranked_response.get("scores", [])
    reranked_documents = [
        texts[idx]
        for idx in sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)
    ]
    print(f"Reranked {len(reranked_documents)} document(s).")
    return reranked_documents

# ---------------------------------------------------------------------------
# LLM chat — Ollama via OpenAI-compatible client
# ---------------------------------------------------------------------------

def llm_chat(
    messages:    list[dict] = [],
    model:       str        = "llama3.1:8b",
    temperature: float      = 0.3,
    max_tokens:  int        = 4096,
) -> str:
    """
    Send a list of messages to the Ollama LLM and return the response as a string.

    Parameters
    ----------
    messages    : list of {"role": str, "content": str} dicts
    model       : Ollama model name (must be pulled via `ollama pull <model>`)
    temperature : creativity control — lower = more deterministic
    max_tokens  : maximum tokens in the response

    Returns
    -------
    str — the assistant's reply, or an error message string on failure
    """
    if model in ("mistral:latest", "phi4:latest"):
        max_tokens = 8192

    try:
        response = ollama_client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=temperature,
            max_tokens=max_tokens,
        )
        content = response.choices[0].message.content
        print(f"LLM response received from model: {model}")
        return content
    except Exception as ex:
        print(f"Error during LLM chat: {ex}")
        return f"An error occurred while contacting Ollama: {ex}"

# ---------------------------------------------------------------------------
# Chat export
# ---------------------------------------------------------------------------

def generate_output_file(chat_history: list[dict]) -> str | None:
    """
    Write the full chat history to a timestamped .txt file in output_dir.
    Returns the file path on success, or None if chat_history is empty.
    """
    if not chat_history:
        print("Chat history is empty — nothing to export.")
        return None

    filename = f"chat_output_{current_datetime}.txt"
    filepath = os.path.join(output_dir, filename)

    with open(filepath, "w", encoding="utf-8") as f:
        for entry in chat_history:
            role    = entry.get("role",    "unknown").capitalize()
            content = entry.get("content", "").strip()
            f.write(f"{role}: {content}\n\n")

    print(f"Chat history exported to: {filepath}")
    return filepath
